"""Report generation for the AI analyst.

Generates Word (.docx), Excel (.xlsx), PDF, and PowerPoint (.pptx) reports
from conversation data. Charts are rendered as images via matplotlib.
"""

import io
import logging
import uuid
from datetime import datetime, timezone

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from openpyxl.chart import BarChart, LineChart, PieChart, Reference
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt as PptxPt
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.colors import HexColor
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage, PageBreak
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.orm import selectinload

from app.db.session import async_session
from app.models.analyst import AnalystConversation, AnalystMessage, AnalystReport, ReportGenStatus
from app.models.notification import Notification, NotificationType
from app.models.user import User
from app.services import email_service, s3

logger = logging.getLogger(__name__)

# Deep Thesis brand colors
BRAND_ACCENT = "#F28C28"
BRAND_ACCENT_HOVER = "#D97A1E"
BRAND_BG = "#FAFAF8"
BRAND_TEXT = "#1A1A1A"
BRAND_TEXT_SECONDARY = "#6B6B6B"
BRAND_TEXT_TERTIARY = "#9B9B9B"
BRAND_BORDER = "#E8E6E3"
BRAND_SCORE_HIGH = "#2D6A4F"
BRAND_SCORE_MID = "#B8860B"
BRAND_SCORE_LOW = "#A23B3B"

CHART_COLORS = [BRAND_ACCENT, BRAND_SCORE_HIGH, BRAND_SCORE_MID, BRAND_SCORE_LOW, "#6366f1", "#ec4899", "#06b6d4", "#84cc16"]


def _render_chart_image(chart_config: dict) -> bytes | None:
    """Render a chart config dict to a PNG image using matplotlib."""
    try:
        chart_type = chart_config.get("type", "bar")
        data = chart_config.get("data", [])
        title = chart_config.get("title", "")
        x_key = chart_config.get("xKey", chart_config.get("nameKey", "name"))
        y_keys = chart_config.get("yKeys", [chart_config.get("dataKey", "value")])
        colors = chart_config.get("colors", CHART_COLORS)

        if not data:
            return None

        fig, ax = plt.subplots(figsize=(8, 5))
        fig.patch.set_facecolor(BRAND_BG)
        ax.set_facecolor(BRAND_BG)
        ax.tick_params(colors=BRAND_TEXT_SECONDARY)
        ax.xaxis.label.set_color(BRAND_TEXT_SECONDARY)
        ax.yaxis.label.set_color(BRAND_TEXT_SECONDARY)
        ax.title.set_color(BRAND_TEXT)
        for spine in ax.spines.values():
            spine.set_color(BRAND_BORDER)

        labels = [str(d.get(x_key, "")) for d in data]

        if chart_type == "pie":
            data_key = y_keys[0] if y_keys else "value"
            values = [d.get(data_key, 0) for d in data]
            ax.pie(values, labels=labels, colors=colors[:len(values)], autopct="%1.1f%%",
                   textprops={"color": BRAND_TEXT})
        elif chart_type == "scatter":
            for i, yk in enumerate(y_keys):
                x_vals = [d.get(x_key, 0) for d in data]
                y_vals = [d.get(yk, 0) for d in data]
                ax.scatter(x_vals, y_vals, color=colors[i % len(colors)], label=yk, alpha=0.7)
            ax.legend(facecolor=BRAND_BG, edgecolor=BRAND_BORDER, labelcolor=BRAND_TEXT_SECONDARY)
        else:
            x = range(len(labels))
            width = 0.8 / len(y_keys) if chart_type == "bar" else 0

            for i, yk in enumerate(y_keys):
                values = [d.get(yk, 0) for d in data]
                color = colors[i % len(colors)]

                if chart_type == "bar":
                    offset = (i - len(y_keys) / 2 + 0.5) * width
                    ax.bar([xi + offset for xi in x], values, width=width, color=color, label=yk)
                elif chart_type == "line":
                    ax.plot(x, values, color=color, label=yk, marker="o", markersize=4)
                elif chart_type == "area":
                    ax.fill_between(x, values, color=color, alpha=0.3, label=yk)
                    ax.plot(x, values, color=color)

            ax.set_xticks(list(x))
            ax.set_xticklabels(labels, rotation=45, ha="right", fontsize=8)
            if len(y_keys) > 1:
                ax.legend(facecolor=BRAND_BG, edgecolor=BRAND_BORDER, labelcolor=BRAND_TEXT_SECONDARY)

        ax.set_title(title, fontsize=12, pad=10)
        plt.tight_layout()

        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=150, facecolor=fig.get_facecolor())
        plt.close(fig)
        buf.seek(0)
        return buf.getvalue()

    except Exception as e:
        logger.warning("Chart rendering failed: %s", e)
        plt.close("all")
        return None


def _generate_docx(conversation: AnalystConversation, messages: list[AnalystMessage], title: str) -> bytes:
    """Generate a Word document from conversation data."""
    doc = Document()

    # Cover page
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run("Deep Thesis Analyst Report")
    run.font.size = Pt(28)
    run.font.color.rgb = RGBColor(0xF2, 0x8C, 0x28)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(title)
    run.font.size = Pt(18)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p.add_run(datetime.now(timezone.utc).strftime("%B %d, %Y"))
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(128, 128, 128)

    doc.add_page_break()

    # Conversation content
    for msg in messages:
        if msg.role == "user" or (hasattr(msg.role, "value") and msg.role.value == "user"):
            doc.add_heading(msg.content[:100], level=2)
        else:
            # Assistant response as body text
            for paragraph_text in msg.content.split("\n\n"):
                paragraph_text = paragraph_text.strip()
                if not paragraph_text:
                    continue
                if paragraph_text.startswith("# "):
                    doc.add_heading(paragraph_text[2:], level=2)
                elif paragraph_text.startswith("## "):
                    doc.add_heading(paragraph_text[3:], level=3)
                elif paragraph_text.startswith("- "):
                    for line in paragraph_text.split("\n"):
                        if line.strip().startswith("- "):
                            doc.add_paragraph(line.strip()[2:], style="List Bullet")
                else:
                    doc.add_paragraph(paragraph_text)

            # Render charts as images
            if msg.charts:
                for chart_config in msg.charts:
                    img_bytes = _render_chart_image(chart_config)
                    if img_bytes:
                        buf = io.BytesIO(img_bytes)
                        doc.add_picture(buf, width=Inches(6))
                        cap = doc.add_paragraph(chart_config.get("title", ""))
                        cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                        cap.style = doc.styles["Caption"] if "Caption" in [s.name for s in doc.styles] else None

    # Citations
    all_citations = []
    for msg in messages:
        if msg.citations:
            all_citations.extend(msg.citations)

    if all_citations:
        doc.add_page_break()
        doc.add_heading("Sources", level=1)
        for i, cite in enumerate(all_citations, 1):
            url = cite.get("url", "") if isinstance(cite, dict) else str(cite)
            title_text = cite.get("title", url) if isinstance(cite, dict) else str(cite)
            doc.add_paragraph(f"{i}. {title_text}\n   {url}", style="List Number")

    # Footer
    section = doc.sections[0]
    footer = section.footer
    footer_para = footer.paragraphs[0] if footer.paragraphs else footer.add_paragraph()
    footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = footer_para.add_run("Generated by Deep Thesis AI Analyst")
    run.font.size = Pt(8)
    run.font.color.rgb = RGBColor(128, 128, 128)

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _generate_xlsx(conversation: AnalystConversation, messages: list[AnalystMessage], title: str) -> bytes:
    """Generate an Excel workbook from conversation data."""
    wb = Workbook()

    # Summary sheet
    ws = wb.active
    ws.title = "Summary"
    ws.append(["Deep Thesis Analyst Report"])
    ws.append([title])
    ws.append([datetime.now(timezone.utc).strftime("%B %d, %Y")])
    ws.append([])
    ws.append(["Question", "Response Summary"])

    # Brand the header row
    header_fill = PatternFill(start_color="F28C28", end_color="F28C28", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True)
    for cell in ws[5]:
        if cell.value:
            cell.fill = header_fill
            cell.font = header_font

    for msg in messages:
        role = msg.role.value if hasattr(msg.role, "value") else msg.role
        if role == "user":
            ws.append([msg.content[:200]])
        elif role == "assistant":
            ws.append(["", msg.content[:500]])

    # Data sheets — one per chart
    chart_num = 0
    for msg in messages:
        if not msg.charts:
            continue
        for chart_config in msg.charts:
            chart_num += 1
            chart_title = chart_config.get("title", f"Chart {chart_num}")
            sheet_name = f"Data {chart_num}"[:31]  # Excel 31-char limit
            ws_data = wb.create_sheet(title=sheet_name)

            data = chart_config.get("data", [])
            if not data:
                continue

            # Headers
            headers = list(data[0].keys())
            ws_data.append(headers)

            # Rows
            for row in data:
                ws_data.append([row.get(h) for h in headers])

            # Add chart
            x_key = chart_config.get("xKey", chart_config.get("nameKey", headers[0]))
            y_keys = chart_config.get("yKeys", [chart_config.get("dataKey")])
            chart_type = chart_config.get("type", "bar")

            try:
                x_col = headers.index(x_key) + 1 if x_key in headers else 1
                chart_obj = None

                if chart_type == "pie":
                    chart_obj = PieChart()
                    data_col = headers.index(y_keys[0]) + 1 if y_keys and y_keys[0] in headers else 2
                    chart_obj.add_data(
                        Reference(ws_data, min_col=data_col, min_row=1, max_row=len(data) + 1),
                        titles_from_data=True,
                    )
                    chart_obj.set_categories(
                        Reference(ws_data, min_col=x_col, min_row=2, max_row=len(data) + 1)
                    )
                elif chart_type in ("line", "area"):
                    chart_obj = LineChart()
                    for yk in y_keys:
                        if yk in headers:
                            col = headers.index(yk) + 1
                            chart_obj.add_data(
                                Reference(ws_data, min_col=col, min_row=1, max_row=len(data) + 1),
                                titles_from_data=True,
                            )
                    chart_obj.set_categories(
                        Reference(ws_data, min_col=x_col, min_row=2, max_row=len(data) + 1)
                    )
                else:
                    chart_obj = BarChart()
                    for yk in y_keys:
                        if yk in headers:
                            col = headers.index(yk) + 1
                            chart_obj.add_data(
                                Reference(ws_data, min_col=col, min_row=1, max_row=len(data) + 1),
                                titles_from_data=True,
                            )
                    chart_obj.set_categories(
                        Reference(ws_data, min_col=x_col, min_row=2, max_row=len(data) + 1)
                    )

                if chart_obj:
                    chart_obj.title = chart_title
                    chart_obj.width = 20
                    chart_obj.height = 12
                    ws_data.add_chart(chart_obj, f"A{len(data) + 4}")
            except Exception as e:
                logger.warning("Excel chart creation failed: %s", e)

    # Sources sheet
    all_citations = []
    for msg in messages:
        if msg.citations:
            all_citations.extend(msg.citations)

    if all_citations:
        ws_sources = wb.create_sheet(title="Sources")
        ws_sources.append(["#", "Title", "URL"])
        for i, cite in enumerate(all_citations, 1):
            url = cite.get("url", "") if isinstance(cite, dict) else str(cite)
            cite_title = cite.get("title", url) if isinstance(cite, dict) else str(cite)
            ws_sources.append([i, cite_title, url])

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue()


def _generate_pdf(conversation: AnalystConversation, messages: list[AnalystMessage], title: str) -> bytes:
    """Generate a PDF document from conversation data."""
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=letter, topMargin=0.75 * inch, bottomMargin=0.75 * inch)
    styles = getSampleStyleSheet()

    # Custom styles
    styles.add(ParagraphStyle(
        "CoverTitle", parent=styles["Title"], fontSize=28,
        textColor=HexColor(BRAND_ACCENT), spaceAfter=12,
    ))
    styles.add(ParagraphStyle(
        "CoverSubtitle", parent=styles["Title"], fontSize=18,
        textColor=HexColor("#333333"), spaceAfter=6,
    ))
    styles.add(ParagraphStyle(
        "CoverDate", parent=styles["Normal"], fontSize=12,
        textColor=HexColor("#808080"), alignment=1, spaceAfter=24,
    ))
    styles.add(ParagraphStyle(
        "UserQuestion", parent=styles["Heading2"], fontSize=14,
        textColor=HexColor(BRAND_ACCENT), spaceBefore=16, spaceAfter=8,
    ))
    styles.add(ParagraphStyle(
        "BodyText2", parent=styles["BodyText"], fontSize=10,
        leading=14, spaceAfter=8,
    ))
    styles.add(ParagraphStyle(
        "BulletItem", parent=styles["BodyText"], fontSize=10,
        leading=14, leftIndent=20, bulletIndent=10, spaceAfter=4,
    ))
    styles.add(ParagraphStyle(
        "SectionH1", parent=styles["Heading2"], fontSize=16,
        textColor=HexColor(BRAND_TEXT), spaceBefore=14, spaceAfter=8,
    ))
    styles.add(ParagraphStyle(
        "SectionH2", parent=styles["Heading3"], fontSize=13,
        textColor=HexColor("#333333"), spaceBefore=10, spaceAfter=6,
    ))
    styles.add(ParagraphStyle(
        "FooterStyle", parent=styles["Normal"], fontSize=8,
        textColor=HexColor("#808080"), alignment=1,
    ))

    story = []

    # Cover page
    story.append(Spacer(1, 2 * inch))
    story.append(Paragraph("Deep Thesis Analyst Report", styles["CoverTitle"]))
    story.append(Paragraph(title, styles["CoverSubtitle"]))
    story.append(Paragraph(datetime.now(timezone.utc).strftime("%B %d, %Y"), styles["CoverDate"]))
    story.append(PageBreak())

    # Conversation content
    for msg in messages:
        role = msg.role.value if hasattr(msg.role, "value") else msg.role
        if role == "user":
            story.append(Paragraph(msg.content[:200], styles["UserQuestion"]))
        else:
            for paragraph_text in msg.content.split("\n\n"):
                paragraph_text = paragraph_text.strip()
                if not paragraph_text:
                    continue
                # Escape XML entities for reportlab
                safe = paragraph_text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                if paragraph_text.startswith("# "):
                    story.append(Paragraph(safe[2:], styles["SectionH1"]))
                elif paragraph_text.startswith("## "):
                    story.append(Paragraph(safe[3:], styles["SectionH2"]))
                elif paragraph_text.startswith("- "):
                    for line in paragraph_text.split("\n"):
                        if line.strip().startswith("- "):
                            safe_line = line.strip()[2:].replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                            story.append(Paragraph(f"• {safe_line}", styles["BulletItem"]))
                else:
                    story.append(Paragraph(safe, styles["BodyText2"]))

            # Charts as images
            if msg.charts:
                for chart_config in msg.charts:
                    img_bytes = _render_chart_image(chart_config)
                    if img_bytes:
                        img_buf = io.BytesIO(img_bytes)
                        story.append(Spacer(1, 12))
                        story.append(RLImage(img_buf, width=6 * inch, height=3.75 * inch))
                        cap_text = chart_config.get("title", "")
                        if cap_text:
                            safe_cap = cap_text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                            story.append(Paragraph(safe_cap, styles["CoverDate"]))
                        story.append(Spacer(1, 12))

    # Citations
    all_citations = []
    for msg in messages:
        if msg.citations:
            all_citations.extend(msg.citations)

    if all_citations:
        story.append(PageBreak())
        story.append(Paragraph("Sources", styles["SectionH1"]))
        for i, cite in enumerate(all_citations, 1):
            url = cite.get("url", "") if isinstance(cite, dict) else str(cite)
            cite_title = cite.get("title", url) if isinstance(cite, dict) else str(cite)
            safe_title = str(cite_title).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            safe_url = str(url).replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(f"{i}. {safe_title}<br/><font size='8' color='#808080'>{safe_url}</font>", styles["BodyText2"]))

    # Footer
    story.append(Spacer(1, 24))
    story.append(Paragraph("Generated by Deep Thesis AI Analyst", styles["FooterStyle"]))

    doc.build(story)
    buf.seek(0)
    return buf.getvalue()


def _generate_pptx(conversation: AnalystConversation, messages: list[AnalystMessage], title: str) -> bytes:
    """Generate a PowerPoint presentation from conversation data."""
    prs = Presentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PptxRGBColor(0xFA, 0xFA, 0xF8)

    txBox = slide.shapes.add_textbox(PptxInches(1), PptxInches(2), PptxInches(11.333), PptxInches(1.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Deep Thesis Analyst Report"
    p.font.size = PptxPt(36)
    p.font.color.rgb = PptxRGBColor(0xF2, 0x8C, 0x28)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = PptxPt(24)
    p2.font.color.rgb = PptxRGBColor(0x1A, 0x1A, 0x1A)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = datetime.now(timezone.utc).strftime("%B %d, %Y")
    p3.font.size = PptxPt(14)
    p3.font.color.rgb = PptxRGBColor(0x6B, 0x6B, 0x6B)
    p3.alignment = PP_ALIGN.CENTER

    # Content slides
    for msg in messages:
        role = msg.role.value if hasattr(msg.role, "value") else msg.role
        if role == "user":
            continue

        # Split long assistant content into slide-sized chunks
        paragraphs = [p.strip() for p in msg.content.split("\n\n") if p.strip()]
        slide_text_chunks = []
        current_chunk = []
        current_len = 0

        for para in paragraphs:
            if current_len + len(para) > 1200 and current_chunk:
                slide_text_chunks.append(current_chunk)
                current_chunk = []
                current_len = 0
            current_chunk.append(para)
            current_len += len(para)
        if current_chunk:
            slide_text_chunks.append(current_chunk)

        for chunk in slide_text_chunks:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            bg = slide.background.fill
            bg.solid()
            bg.fore_color.rgb = PptxRGBColor(0xFA, 0xFA, 0xF8)

            txBox = slide.shapes.add_textbox(PptxInches(0.75), PptxInches(0.5), PptxInches(11.833), PptxInches(6.5))
            tf = txBox.text_frame
            tf.word_wrap = True

            for i, para in enumerate(chunk):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                if para.startswith("# "):
                    p.text = para[2:]
                    p.font.size = PptxPt(24)
                    p.font.color.rgb = PptxRGBColor(0xF2, 0x8C, 0x28)
                    p.font.bold = True
                elif para.startswith("## "):
                    p.text = para[3:]
                    p.font.size = PptxPt(20)
                    p.font.color.rgb = PptxRGBColor(0x1A, 0x1A, 0x1A)
                    p.font.bold = True
                elif para.startswith("- "):
                    for line in para.split("\n"):
                        if line.strip().startswith("- "):
                            bp = tf.add_paragraph() if p.text else p
                            bp.text = f"  •  {line.strip()[2:]}"
                            bp.font.size = PptxPt(14)
                            bp.font.color.rgb = PptxRGBColor(0x6B, 0x6B, 0x6B)
                else:
                    p.text = para[:500]
                    p.font.size = PptxPt(14)
                    p.font.color.rgb = PptxRGBColor(0x6B, 0x6B, 0x6B)

        # Chart slides
        if msg.charts:
            for chart_config in msg.charts:
                img_bytes = _render_chart_image(chart_config)
                if img_bytes:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    bg = slide.background.fill
                    bg.solid()
                    bg.fore_color.rgb = PptxRGBColor(0xFA, 0xFA, 0xF8)

                    # Chart title
                    chart_title = chart_config.get("title", "")
                    if chart_title:
                        txBox = slide.shapes.add_textbox(PptxInches(0.75), PptxInches(0.3), PptxInches(11.833), PptxInches(0.6))
                        tf = txBox.text_frame
                        p = tf.paragraphs[0]
                        p.text = chart_title
                        p.font.size = PptxPt(20)
                        p.font.color.rgb = PptxRGBColor(0x1A, 0x1A, 0x1A)
                        p.alignment = PP_ALIGN.CENTER

                    img_buf = io.BytesIO(img_bytes)
                    slide.shapes.add_picture(img_buf, PptxInches(1.5), PptxInches(1.2), PptxInches(10.333), PptxInches(5.8))

    # Sources slide
    all_citations = []
    for msg in messages:
        if msg.citations:
            all_citations.extend(msg.citations)

    if all_citations:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = PptxRGBColor(0xFA, 0xFA, 0xF8)

        txBox = slide.shapes.add_textbox(PptxInches(0.75), PptxInches(0.5), PptxInches(11.833), PptxInches(6.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Sources"
        p.font.size = PptxPt(24)
        p.font.color.rgb = PptxRGBColor(0xF2, 0x8C, 0x28)
        p.font.bold = True

        for i, cite in enumerate(all_citations[:20], 1):
            url = cite.get("url", "") if isinstance(cite, dict) else str(cite)
            cite_title = cite.get("title", url) if isinstance(cite, dict) else str(cite)
            p = tf.add_paragraph()
            p.text = f"{i}. {cite_title}"
            p.font.size = PptxPt(11)
            p.font.color.rgb = PptxRGBColor(0x6B, 0x6B, 0x6B)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


async def generate_report(report_id: str) -> None:
    """Generate a report (background task). Updates status in DB and uploads to S3."""
    rid = uuid.UUID(report_id)

    async with async_session() as db:
        try:
            # Load report with conversation and messages
            result = await db.execute(
                select(AnalystReport).where(AnalystReport.id == rid)
            )
            report = result.scalar_one_or_none()
            if not report:
                logger.error("Report %s not found", report_id)
                return

            report.status = ReportGenStatus.generating.value
            await db.commit()

            # Load conversation with messages
            result = await db.execute(
                select(AnalystConversation)
                .where(AnalystConversation.id == report.conversation_id)
                .options(selectinload(AnalystConversation.messages))
            )
            conversation = result.scalar_one()
            all_messages = list(conversation.messages)
            # Use only the last assistant message for the report
            assistant_messages = [
                m for m in all_messages
                if (m.role.value if hasattr(m.role, "value") else m.role) == "assistant"
            ]
            messages = [assistant_messages[-1]] if assistant_messages else all_messages

            # Generate document
            fmt = report.format.value if hasattr(report.format, "value") else report.format
            if fmt == "docx":
                file_bytes = _generate_docx(conversation, messages, report.title)
                ext = "docx"
                content_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            elif fmt == "pdf":
                file_bytes = _generate_pdf(conversation, messages, report.title)
                ext = "pdf"
                content_type = "application/pdf"
            elif fmt == "pptx":
                file_bytes = _generate_pptx(conversation, messages, report.title)
                ext = "pptx"
                content_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            else:
                file_bytes = _generate_xlsx(conversation, messages, report.title)
                ext = "xlsx"
                content_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

            # Upload to S3
            s3_key = f"analyst-reports/{conversation.id}/{report.id}/report.{ext}"
            s3.upload_file(file_bytes, s3_key)

            report.s3_key = s3_key
            report.file_size_bytes = len(file_bytes)
            report.status = ReportGenStatus.complete.value
            await db.commit()

            # Create notification for user
            fmt_label = ext.upper()
            notification = Notification(
                user_id=report.user_id,
                type=NotificationType.report_ready,
                title="Report ready",
                message=f"{fmt_label} report",
                link=f"/api/analyst/reports/{report.id}/download",
            )
            db.add(notification)
            await db.commit()

            # Send email notification
            user_result = await db.execute(select(User).where(User.id == report.user_id))
            user = user_result.scalar_one_or_none()
            if user:
                email_service.send_report_ready(
                    user_email=user.email,
                    user_name=user.name,
                    report_format=ext,
                )

            logger.info("Report %s generated: %s (%d bytes)", report_id, s3_key, len(file_bytes))

        except Exception as e:
            logger.error("Report generation failed for %s: %s", report_id, e)
            try:
                report.status = ReportGenStatus.failed.value
                report.error = str(e)[:500]
                await db.commit()
            except Exception:
                logger.error("Failed to update report status for %s", report_id)
