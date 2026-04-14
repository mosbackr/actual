import csv
import io
import subprocess
import tempfile
from pathlib import Path


def extract_text(file_data: bytes, filename: str, file_type: str) -> str:
    extractors = {
        "pdf": _extract_pdf,
        "docx": _extract_docx,
        "doc": _extract_doc,
        "pptx": _extract_pptx,
        "ppt": _extract_ppt,
        "xlsx": _extract_xlsx,
        "xls": _extract_xls,
        "csv": _extract_csv,
        "md": _extract_text,
        "txt": _extract_text,
    }
    extractor = extractors.get(file_type)
    if not extractor:
        return f"[Unsupported file type: {file_type}]"
    try:
        return extractor(file_data, filename)
    except Exception as e:
        return f"[Error extracting {filename}: {e}]"


def _extract_pdf(data: bytes, filename: str) -> str:
    import fitz
    doc = fitz.open(stream=data, filetype="pdf")
    pages = []
    for i, page in enumerate(doc):
        text = page.get_text()
        if text.strip():
            pages.append(f"--- Page {i + 1} ---\n{text.strip()}")
    doc.close()
    return "\n\n".join(pages) if pages else "[No readable text found in PDF]"


def _extract_docx(data: bytes, filename: str) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs) if paragraphs else "[Empty document]"


def _extract_doc(data: bytes, filename: str) -> str:
    with tempfile.TemporaryDirectory() as tmpdir:
        doc_path = Path(tmpdir) / filename
        doc_path.write_bytes(data)
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "docx", "--outdir", tmpdir, str(doc_path)],
            capture_output=True, timeout=30,
        )
        docx_path = doc_path.with_suffix(".docx")
        if docx_path.exists():
            return _extract_docx(docx_path.read_bytes(), docx_path.name)
    return "[Failed to convert .doc file]"


def _extract_pptx(data: bytes, filename: str) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append(f"--- Slide {i + 1} ---\n" + "\n".join(texts))
    return "\n\n".join(slides) if slides else "[Empty presentation]"


def _extract_ppt(data: bytes, filename: str) -> str:
    with tempfile.TemporaryDirectory() as tmpdir:
        ppt_path = Path(tmpdir) / filename
        ppt_path.write_bytes(data)
        subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "pptx", "--outdir", tmpdir, str(ppt_path)],
            capture_output=True, timeout=30,
        )
        pptx_path = ppt_path.with_suffix(".pptx")
        if pptx_path.exists():
            return _extract_pptx(pptx_path.read_bytes(), pptx_path.name)
    return "[Failed to convert .ppt file]"


def _extract_xlsx(data: bytes, filename: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheets = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
        if rows:
            header = rows[0]
            sep = " | ".join(["---"] * len(header.split(" | ")))
            table = f"### Sheet: {sheet_name}\n\n{header}\n{sep}\n" + "\n".join(rows[1:])
            sheets.append(table)
    wb.close()
    return "\n\n".join(sheets) if sheets else "[Empty spreadsheet]"


def _extract_xls(data: bytes, filename: str) -> str:
    import xlrd
    wb = xlrd.open_workbook(file_contents=data)
    sheets = []
    for sheet in wb.sheets():
        rows = []
        for row_idx in range(sheet.nrows):
            cells = [str(sheet.cell_value(row_idx, col)) for col in range(sheet.ncols)]
            if any(c.strip() for c in cells):
                rows.append(" | ".join(cells))
        if rows:
            header = rows[0]
            sep = " | ".join(["---"] * len(header.split(" | ")))
            table = f"### Sheet: {sheet.name}\n\n{header}\n{sep}\n" + "\n".join(rows[1:])
            sheets.append(table)
    return "\n\n".join(sheets) if sheets else "[Empty spreadsheet]"


def _extract_csv(data: bytes, filename: str) -> str:
    text = data.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = []
    for row in reader:
        if any(c.strip() for c in row):
            rows.append(" | ".join(row))
    if rows:
        header = rows[0]
        sep = " | ".join(["---"] * len(header.split(" | ")))
        return f"{header}\n{sep}\n" + "\n".join(rows[1:])
    return "[Empty CSV]"


def _extract_text(data: bytes, filename: str) -> str:
    return data.decode("utf-8", errors="replace")


def consolidate_documents(documents: list[dict]) -> str:
    type_labels = {
        "pdf": "document", "docx": "document", "doc": "document",
        "pptx": "slides", "ppt": "slides",
        "xlsx": "spreadsheet", "xls": "spreadsheet", "csv": "spreadsheet",
        "md": "markdown", "txt": "text",
    }
    sections = []
    for doc in documents:
        label = type_labels.get(doc["file_type"], "file")
        sections.append(f"=== DOCUMENT: {doc['filename']} ({label}) ===\n\n{doc['text']}")
    return "\n\n".join(sections)
